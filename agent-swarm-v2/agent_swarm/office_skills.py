"""Office Skills — generate .xlsx, .pptx, .docx files as agent tools.

Inspired by Aden Hive's native Office Suite Skills. Agents can produce
professional office artifacts without manual post-processing.
"""

import logging
import os

logger = logging.getLogger(__name__)


def generate_xlsx(filename: str, data: dict, workspace: str = "") -> str:
    """Generate an Excel file from structured data.

    Args:
        filename: output filename (e.g. 'report.xlsx')
        data: dict with 'sheets' key, each sheet is a list of dict rows
              e.g. {"sheets": {"Sheet1": [{"Name": "Alice", "Score": 95}, ...]}}
        workspace: working directory to save file
    """
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        return "Error: openpyxl not installed. Install with: pip install openpyxl"

    try:
        wb = Workbook()
        wb.remove(wb.active)  # remove default sheet

        sheets_data = data.get("sheets", {"Sheet1": data.get("rows", [])})
        header_font = Font(bold=True, size=11)
        header_fill = PatternFill(start_color="E2E8F0", end_color="E2E8F0", fill_type="solid")
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        for sheet_name, rows in sheets_data.items():
            ws = wb.create_sheet(title=sheet_name[:31])  # Excel 31-char limit
            if not rows:
                ws.cell(row=1, column=1, value="(empty)")
                continue

            # Write headers
            headers = list(rows[0].keys()) if isinstance(rows[0], dict) else [f"Col{i}" for i in range(len(rows[0]))]
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border

            # Write data
            for row_idx, row in enumerate(rows, 2):
                if isinstance(row, dict):
                    for col_idx, header in enumerate(headers, 1):
                        cell = ws.cell(row=row_idx, column=col_idx, value=row.get(header, ""))
                        cell.border = thin_border
                else:
                    for col_idx, val in enumerate(row, 1):
                        cell = ws.cell(row=row_idx, column=col_idx, value=val)
                        cell.border = thin_border

            # Auto-fit column widths
            for col_idx, header in enumerate(headers, 1):
                max_width = len(str(header)) + 2
                for row_idx in range(2, len(rows) + 2):
                    val = ws.cell(row=row_idx, column=col_idx).value
                    if val:
                        max_width = max(max_width, min(len(str(val)), 40) + 2)
                ws.column_dimensions[ws.cell(row=1, column=col_idx).column_letter].width = max_width

        save_path = os.path.join(workspace, filename) if workspace else filename
        os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
        wb.save(save_path)
        total_rows = sum(len(rows) for rows in sheets_data.values())
        return f"Excel file saved: {save_path} ({len(sheets_data)} sheets, {total_rows} rows)"
    except Exception as e:
        return f"Error generating Excel: {e}"


def generate_docx(filename: str, content: dict, workspace: str = "") -> str:
    """Generate a Word document from structured data.

    Args:
        filename: output filename (e.g. 'report.docx')
        content: dict with 'title', 'sections' list (each with 'heading' and 'body')
        workspace: working directory
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches
    except ImportError:
        return "Error: python-docx not installed. Install with: pip install python-docx"

    try:
        doc = Document()
        # Title
        title = content.get("title", "Report")
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = 1  # center

        # Metadata
        if content.get("subtitle"):
            doc.add_paragraph(content["subtitle"], style="Subtitle")

        # Sections
        for section in content.get("sections", []):
            heading = section.get("heading", "")
            body = section.get("body", "")
            if heading:
                doc.add_heading(heading, level=1)
            for para_text in body.split("\n"):
                para_text = para_text.strip()
                if para_text:
                    doc.add_paragraph(para_text)

        save_path = os.path.join(workspace, filename) if workspace else filename
        os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
        doc.save(save_path)
        return f"Word document saved: {save_path} ({len(content.get('sections', []))} sections)"
    except Exception as e:
        return f"Error generating Word document: {e}"


def generate_pptx(filename: str, content: dict, workspace: str = "") -> str:
    """Generate a PowerPoint presentation from structured data.

    Args:
        filename: output filename (e.g. 'presentation.pptx')
        content: dict with 'title', 'slides' list (each with 'title' and 'bullets' list)
        workspace: working directory
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt, Inches
    except ImportError:
        return "Error: python-pptx not installed. Install with: pip install python-pptx"

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = content.get("title", "Presentation")
        if content.get("subtitle"):
            slide.placeholders[1].text = content["subtitle"]

        # Content slides
        for slide_data in content.get("slides", []):
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data.get("title", "")

            bullets = slide_data.get("bullets", [])
            if bullets:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        tf.text = bullet
                    else:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 0

        save_path = os.path.join(workspace, filename) if workspace else filename
        os.makedirs(os.path.dirname(save_path) or ".", exist_ok=True)
        prs.save(save_path)
        return f"PowerPoint saved: {save_path} ({len(content.get('slides', []))} slides)"
    except Exception as e:
        return f"Error generating PowerPoint: {e}"
